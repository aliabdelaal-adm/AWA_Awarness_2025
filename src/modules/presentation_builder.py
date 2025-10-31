"""
Presentation Builder Module
Creates professional PowerPoint presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from arabic_reshaper import reshape
from bidi.algorithm import get_display
import os
from typing import List, Dict, Optional


class PresentationBuilder:
    """Build professional PowerPoint presentations"""
    
    def __init__(self):
        """Initialize presentation builder"""
        self.presentation = None
    
    def _format_arabic_text(self, text: str) -> str:
        """
        Format Arabic text for proper display in PowerPoint
        
        Args:
            text: Arabic text to format
            
        Returns:
            Properly formatted Arabic text
        """
        try:
            reshaped_text = reshape(text)
            bidi_text = get_display(reshaped_text)
            return bidi_text
        except:
            # If reshaping fails, return original text
            return text
        
    def create_presentation(self, title: str, slides: List[Dict], 
                          output_path: str, language: str = 'ar', 
                          max_slides: int = None) -> str:
        """
        Create a PowerPoint presentation
        
        Args:
            title: Presentation title
            slides: List of slide dictionaries with 'title' and 'content' keys
            output_path: Path to save the presentation
            language: Language code ('ar' or 'en')
            max_slides: Maximum number of slides to create (None for all)
            
        Returns:
            Path to created presentation file
        """
        try:
            # Create presentation
            self.presentation = Presentation()
            
            # Set slide size to 16:9
            self.presentation.slide_width = Inches(13.333)
            self.presentation.slide_height = Inches(7.5)
            
            # Create title slide
            self._create_title_slide(title, language)
            
            # Limit slides if max_slides specified
            slides_to_create = slides
            if max_slides and max_slides > 0:
                # Smart merging: combine slides to fit within max_slides
                slides_to_create = self._merge_slides_intelligently(slides, max_slides)
            
            # Create content slides
            for slide_data in slides_to_create:
                if 'image_path' in slide_data:
                    self._create_image_slide(
                        slide_data.get('title', ''),
                        slide_data.get('image_path', ''),
                        slide_data.get('content', ''),
                        language
                    )
                else:
                    self._create_content_slide(
                        slide_data.get('title', ''),
                        slide_data.get('content', ''),
                        language
                    )
            
            # Save presentation
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.presentation.save(output_path)
            
            print(f"✓ PowerPoint presentation created: {output_path}")
            print(f"  Total slides: {len(self.presentation.slides)}")
            return output_path
            
        except Exception as e:
            print(f"Error creating presentation: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _merge_slides_intelligently(self, slides: List[Dict], max_slides: int) -> List[Dict]:
        """
        Intelligently merge slides to reduce total count
        
        Args:
            slides: Original list of slides
            max_slides: Maximum number of slides desired
            
        Returns:
            Merged list of slides
        """
        if len(slides) <= max_slides:
            return slides
        
        # Calculate how many slides to merge together
        merge_ratio = len(slides) / max_slides
        merged_slides = []
        
        i = 0
        while i < len(slides) and len(merged_slides) < max_slides:
            # Determine how many slides to merge
            slides_to_merge = int(merge_ratio) if len(merged_slides) < max_slides - 1 else len(slides) - i
            slides_to_merge = max(1, min(slides_to_merge, len(slides) - i))
            
            # Merge slides
            if slides_to_merge == 1:
                merged_slides.append(slides[i])
            else:
                # Combine multiple slides into one
                merged_slide = {
                    'title': slides[i].get('title', f'Section {len(merged_slides) + 1}'),
                    'content': '',
                    'source': 'merged'
                }
                
                # Combine content from multiple slides
                content_parts = []
                for j in range(i, min(i + slides_to_merge, len(slides))):
                    slide = slides[j]
                    if 'content' in slide and slide['content']:
                        content_parts.append(slide['content'])
                    elif 'image_path' in slide:
                        # Keep image slides separate
                        merged_slides.append(slide)
                        continue
                
                # Join content with bullet points
                merged_slide['content'] = '\n\n'.join(content_parts)
                if content_parts:
                    merged_slides.append(merged_slide)
            
            i += slides_to_merge
        
        return merged_slides
    
    def _create_title_slide(self, title: str, language: str = 'ar'):
        """Create title slide"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Format title for Arabic if needed
        display_title = self._format_arabic_text(title) if language == 'ar' else title
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = display_title
        
        # Format title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Set font for Arabic
        if language == 'ar':
            title_frame.paragraphs[0].font.name = 'Arial'
        
        # Add subtitle if Arabic
        if language == 'ar' and slide.placeholders:
            try:
                subtitle = slide.placeholders[1]
                subtitle_text = self._format_arabic_text("عرض إحترافي مبتكر")
                subtitle.text = subtitle_text
                subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                subtitle.text_frame.paragraphs[0].font.size = Pt(32)
                subtitle.text_frame.paragraphs[0].font.name = 'Arial'
            except:
                pass
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue
    
    def _create_content_slide(self, slide_title: str, content: str, language: str = 'ar'):
        """Create content slide with title and text"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Format title for Arabic if needed
        display_title = self._format_arabic_text(slide_title) if language == 'ar' else slide_title
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = display_title
        
        # Format title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        if language == 'ar':
            title_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            title_frame.paragraphs[0].font.name = 'Arial'
        
        # Add content
        if slide.placeholders:
            try:
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                # Format content for Arabic if needed
                display_content = self._format_arabic_text(content) if language == 'ar' else content
                
                # Split content into paragraphs
                paragraphs = display_content.split('\n')
                
                for i, para_text in enumerate(paragraphs):
                    if para_text.strip():
                        if i > 0:
                            text_frame.add_paragraph()
                        
                        p = text_frame.paragraphs[-1] if i > 0 else text_frame.paragraphs[0]
                        p.text = para_text.strip()
                        p.font.size = Pt(24)
                        p.level = 0
                        
                        if language == 'ar':
                            p.alignment = PP_ALIGN.RIGHT
                            p.font.name = 'Arial'
                        else:
                            p.alignment = PP_ALIGN.LEFT
            except Exception as e:
                print(f"Error adding content to slide: {e}")
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    
    def _create_image_slide(self, slide_title: str, image_path: str, caption: str = '', language: str = 'ar'):
        """Create slide with image and text merged professionally"""
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Format title for Arabic if needed
        display_title = self._format_arabic_text(slide_title) if language == 'ar' else slide_title
        
        # Add title at top
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12.333)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = display_title
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.alignment = PP_ALIGN.CENTER
        if language == 'ar':
            p.font.name = 'Arial'
        
        # Add image
        if os.path.exists(image_path):
            try:
                # Calculate image position and size - leave space for caption
                img_left = Inches(1.5)
                img_top = Inches(2)
                img_width = Inches(10.333)
                img_height = Inches(4)  # Reduced to leave space for caption
                
                slide.shapes.add_picture(
                    image_path,
                    img_left,
                    img_top,
                    width=img_width,
                    height=img_height
                )
            except Exception as e:
                print(f"Error adding image to slide: {e}")
        
        # Add caption if provided - merged with image
        if caption:
            display_caption = self._format_arabic_text(caption) if language == 'ar' else caption
            
            caption_left = Inches(0.5)
            caption_top = Inches(6.2)
            caption_width = Inches(12.333)
            caption_height = Inches(1)
            
            caption_box = slide.shapes.add_textbox(
                caption_left, caption_top, caption_width, caption_height
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = display_caption
            caption_frame.word_wrap = True
            
            p = caption_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
            if language == 'ar':
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.RIGHT
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    
    def add_custom_template(self, template_path: str) -> bool:
        """
        Load a custom PowerPoint template
        
        Args:
            template_path: Path to .pptx template file
            
        Returns:
            True if successful, False otherwise
        """
        try:
            if os.path.exists(template_path):
                self.presentation = Presentation(template_path)
                return True
            else:
                print(f"Template not found: {template_path}")
                return False
        except Exception as e:
            print(f"Error loading template: {e}")
            return False
