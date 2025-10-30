"""
PowerPoint Presentation Generator Module
Creates professional PowerPoint presentations from PDF content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional
import os


class PresentationGenerator:
    """Generate professional PowerPoint presentations"""
    
    def __init__(self, config: dict = None):
        """
        Initialize presentation generator
        
        Args:
            config: Configuration dictionary
        """
        self.config = config or self.get_default_config()
        self.prs = None
        
    def get_default_config(self) -> dict:
        """Get default configuration"""
        return {
            'theme': 'professional',
            'title_font_size': 44,
            'content_font_size': 18,
            'background_color': (255, 255, 255),
            'title_color': (0, 51, 102),
            'content_color': (0, 0, 0),
            'slide_layout': 'title_and_content'
        }
    
    def create_presentation(self, title: str = "Presentation") -> Presentation:
        """
        Create a new presentation
        
        Args:
            title: Presentation title
            
        Returns:
            Presentation object
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        return self.prs
    
    def add_title_slide(self, title: str, subtitle: str = None):
        """
        Add title slide
        
        Args:
            title: Main title
            subtitle: Subtitle text
        """
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(self.config.get('title_font_size', 44))
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.config.get('title_color', (0, 51, 102)))
        
        # Set subtitle if provided
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(28)
    
    def add_content_slide(self, title: str, content: str):
        """
        Add content slide with title and body
        
        Args:
            title: Slide title
            content: Slide content
        """
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.config.get('title_color', (0, 51, 102)))
        
        # Set content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add content
        p = text_frame.paragraphs[0]
        p.text = content
        p.font.size = Pt(self.config.get('content_font_size', 18))
        p.font.color.rgb = RGBColor(*self.config.get('content_color', (0, 0, 0)))
        p.level = 0
    
    def add_bullet_slide(self, title: str, bullets: List[str]):
        """
        Add slide with bullet points
        
        Args:
            title: Slide title
            bullets: List of bullet points
        """
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.config.get('title_color', (0, 51, 102)))
        
        # Add bullets
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = Pt(self.config.get('content_font_size', 18))
            p.font.color.rgb = RGBColor(*self.config.get('content_color', (0, 0, 0)))
    
    def generate_from_chunks(self, text_chunks: List[str], title: str, 
                            subtitle: str = None, output_path: str = None) -> str:
        """
        Generate presentation from text chunks
        
        Args:
            text_chunks: List of text chunks (each becomes a slide)
            title: Presentation title
            subtitle: Presentation subtitle
            output_path: Output file path
            
        Returns:
            Path to generated presentation
        """
        # Create presentation
        self.create_presentation(title)
        
        # Add title slide
        self.add_title_slide(title, subtitle)
        
        # Add content slides
        for i, chunk in enumerate(text_chunks, 1):
            # Extract title from chunk or use default
            lines = chunk.strip().split('\n')
            if lines:
                slide_title = lines[0][:100] if len(lines[0]) > 100 else lines[0]
                slide_content = '\n'.join(lines[1:]) if len(lines) > 1 else chunk
            else:
                slide_title = f"Slide {i}"
                slide_content = chunk
            
            # Check if content has bullet-like format
            if any(line.strip().startswith(('-', '•', '*', '○')) for line in lines):
                bullets = [line.strip().lstrip('-•*○ ') for line in lines if line.strip().startswith(('-', '•', '*', '○'))]
                self.add_bullet_slide(slide_title, bullets)
            else:
                self.add_content_slide(slide_title, slide_content)
        
        # Save presentation
        if output_path is None:
            output_dir = os.path.join('output', 'presentations')
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{title.replace(' ', '_')}.pptx")
        
        self.prs.save(output_path)
        return output_path
    
    def add_closing_slide(self, message: str = "Thank You"):
        """
        Add closing/thank you slide
        
        Args:
            message: Closing message
        """
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add text box with centered message
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(1.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = message
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.config.get('title_color', (0, 51, 102)))
